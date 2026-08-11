from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .config import COLORS, FONT, LOGO_PATH, SAFE_X, SAFE_Y, SLIDE_H, SLIDE_W, SOURCE_TEXT


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.replace("#", "")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def prs_new() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: str = "white") -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(COLORS[color])


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 14, color: str = "dark_text", bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(COLORS[color])
    return box


def add_title(slide, title: str, subtitle: str | None = None, x: float = SAFE_X, y: float = 0.42, w: float = 8.5) -> None:
    add_text(slide, title, x, y, w, 0.55, size=30, color="navy", bold=True)
    if subtitle:
        add_text(slide, subtitle, x, y + 0.58, w, 0.45, size=13, color="dark_text")
    add_accent(slide, x, y + 1.1, 0.85, "cyan")


def add_accent(slide, x: float, y: float, w: float = 0.7, color: str = "cyan") -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color])
    shape.line.fill.background()


def add_footer(slide, text: str = SOURCE_TEXT) -> None:
    add_text(slide, text, SAFE_X, 7.16, 8.8, 0.18, size=7, color="mid_gray")
    if LOGO_PATH.exists():
        try:
            slide.shapes.add_picture(str(LOGO_PATH), Inches(11.7), Inches(6.88), width=Inches(1.05))
        except Exception:
            pass


def add_card(slide, title: str, value: str, subtitle: str, x: float, y: float, w: float, h: float, color: str = "cyan", text_color: str = "white"):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.04), Inches(y + 0.05), Inches(w), Inches(h))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = rgb("E6E9F0")
    shadow.line.fill.background()
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color])
    shape.line.fill.background()
    add_text(slide, title, x + 0.15, y + 0.14, w - 0.3, 0.22, size=9, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, value, x + 0.12, y + 0.38, w - 0.24, 0.5, size=27, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, x + 0.15, y + 0.93, w - 0.3, 0.25, size=8, color=text_color, align=PP_ALIGN.CENTER)
    return shape


def add_light_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, accent: str = "cyan"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb(COLORS["light_gray"])
    add_accent(slide, x + 0.16, y + 0.18, 0.44, accent)
    add_text(slide, title, x + 0.18, y + 0.34, w - 0.36, 0.28, size=12, color="navy", bold=True)
    add_text(slide, body, x + 0.18, y + 0.72, w - 0.36, h - 0.82, size=10, color="dark_text")
    return shape


def add_big_number(slide, value: str, label: str, x: float, y: float, color: str = "cyan"):
    add_text(slide, value, x, y, 2.0, 0.6, size=34, color=color, bold=True)
    add_text(slide, label, x, y + 0.62, 2.5, 0.35, size=10, color="dark_text")


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size: int = 13, color: str = "dark_text"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(COLORS[color])
        p.space_after = Pt(8)
    return box


def add_image(slide, image_path: Path, x: float, y: float, w: float, h: float | None = None):
    if h is None:
        return slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, header_color: str = "navy"):
    if not rows:
        return None
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = rgb(COLORS[header_color] if r == 0 else ("F8FAFC" if r % 2 else "FFFFFF"))
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(9 if r else 8)
                p.font.bold = r == 0
                p.font.color.rgb = rgb("FFFFFF" if r == 0 else COLORS["dark_text"])
    return table_shape


def add_section_label(slide, text: str, x: float, y: float, color: str = "magenta"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.35), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[color])
    shape.line.fill.background()
    add_text(slide, text, x + 0.08, y + 0.07, 1.18, 0.18, size=8, color="white", bold=True, align=PP_ALIGN.CENTER)


def add_note(slide, text: str, x: float, y: float, w: float, h: float):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("F7F9FC")
    shape.line.color.rgb = rgb(COLORS["light_gray"])
    add_text(slide, text, x + 0.18, y + 0.16, w - 0.36, h - 0.24, size=11, color="dark_text")
